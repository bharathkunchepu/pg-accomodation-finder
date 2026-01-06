"""
Script to generate PowerPoint presentation from the outline
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def create_presentation():
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    purple = RGBColor(124, 58, 237)  # #7c3aed
    indigo = RGBColor(99, 102, 241)  # #6366f1
    gray = RGBColor(107, 114, 128)   # #6b7280
    dark_gray = RGBColor(31, 41, 55)  # #1f2937
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "PG Accommodation Finder"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(left, Inches(3.5), width, Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Complete Frontend React Application"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = gray
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    tagline_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Find Your Perfect Paying Guest Accommodation"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(20)
    tagline_para.font.color.rgb = dark_gray
    tagline_para.alignment = PP_ALIGN.CENTER
    
    footer_box = slide.shapes.add_textbox(left, Inches(6), width, Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Built with React, React Router, and TailwindCSS"
    footer_para = footer_frame.paragraphs[0]
    footer_para.font.size = Pt(14)
    footer_para.font.color.rgb = gray
    footer_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Problem Statement"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Students and professionals struggle to find suitable PG accommodations"
    p = tf.add_paragraph()
    p.text = "Lack of centralized platform for PG listings"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Difficulty in filtering by budget, location, and amenities"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "No easy way for owners to manage their listings"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Limited access to reviews and ratings"
    p.level = 0
    
    # Solution box
    solution_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    solution_frame = solution_box.text_frame
    solution_frame.text = "Solution: A comprehensive web application that connects students with PG owners, providing search, filter, booking, and management capabilities."
    solution_para = solution_frame.paragraphs[0]
    solution_para.font.size = Pt(14)
    solution_para.font.bold = True
    solution_para.font.color.rgb = RGBColor(0, 0, 0)
    solution_frame.margin_left = Inches(0.2)
    solution_frame.margin_right = Inches(0.2)
    solution_frame.margin_top = Inches(0.1)
    solution_frame.margin_bottom = Inches(0.1)
    fill = solution_box.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(254, 243, 199)  # Light yellow
    
    # Slide 3: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "A full-featured frontend-only React application for finding and managing Paying Guest (PG) accommodations."
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Statistics:"
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• 10+ Pages"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 2 User Types (Students & Owners)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• 100% Frontend"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Points:"
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Complete CRUD operations for PG listings"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Advanced search and filtering system"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Booking management system"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Review and rating functionality"
    p.level = 1
    
    # Slide 4: Key Features - For Users
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features - For Users"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🔍 Advanced Search"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • Search by name, location, area"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Real-time filtering"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Multiple filter combinations"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "💰 Budget Filtering"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • Min/Max price range"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Sort by price (low to high)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Affordable options highlighted"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "⭐ Reviews & Ratings"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • View authentic reviews"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Add your own reviews"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Star rating system (1-5)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "📋 Booking System"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • Request bookings online"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Track booking status"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Booking confirmation"
    p.level = 1
    
    # Slide 5: Key Features - For Owners
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features - For Owners"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "📊 Dashboard"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • View all listings"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Booking statistics"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Quick status overview"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "➕ Add/Edit PGs"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • Create new listings"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Edit existing PGs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Upload multiple images"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "✅ Booking Management"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • View booking requests"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Approve/Reject bookings"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Track booking status"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "🔄 Status Control"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "  • Mark as Vacant/Occupied"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Update room availability"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "  • Real-time status updates"
    p.level = 1
    
    # Slide 6: Technology Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Technologies:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "• React 18"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• React Router DOM"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Context API"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• TailwindCSS"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Vite"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• LocalStorage"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Why These Technologies?"
    p.font.bold = True
    p.font.size = Pt(20)
    p = tf.add_paragraph()
    p.text = "• React: Component-based architecture, reusable UI"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• React Router: Client-side routing, seamless navigation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Context API: Global state management without Redux"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• TailwindCSS: Utility-first CSS, rapid UI development"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Vite: Fast build tool, excellent dev experience"
    p.level = 1
    
    # Slide 7: Project Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Project Structure:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "pg-accommodation-finder/"
    p.font.name = "Courier New"
    p.font.size = Pt(12)
    p = tf.add_paragraph()
    p.text = "├── public/pgs.json"
    p.font.name = "Courier New"
    p.font.size = Pt(11)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "├── src/components/ (Navbar, Footer, PGCard, etc.)"
    p.font.name = "Courier New"
    p.font.size = Pt(11)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "├── src/context/PGContext.jsx"
    p.font.name = "Courier New"
    p.font.size = Pt(11)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "├── src/pages/ (Home, Listings, Dashboard, etc.)"
    p.font.name = "Courier New"
    p.font.size = Pt(11)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "└── src/App.jsx"
    p.font.name = "Courier New"
    p.font.size = Pt(11)
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Point: Clean Architecture - Separation of concerns with components, pages, and context"
    p.font.italic = True
    p.font.size = Pt(14)
    
    # Slide 7.5: System Architecture Diagram
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture & Prototype"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add the diagram image if it exists
    import os
    diagram_path = "PG_Finder_Architecture_Diagram.png"
    if os.path.exists(diagram_path):
        left_img = Inches(0.5)
        top_img = Inches(1.2)
        width_img = Inches(9)
        height_img = Inches(6)
        slide.shapes.add_picture(diagram_path, left_img, top_img, width_img, height_img)
    else:
        # Fallback text if image not found
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
        desc_frame = desc_box.text_frame
        desc_frame.text = "Architecture Diagram: PG_Finder_Architecture_Diagram.png"
        desc_para = desc_frame.paragraphs[0]
        desc_para.font.size = Pt(14)
        desc_para.alignment = PP_ALIGN.CENTER
    
    # Slide 8: Application Pages
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Application Pages"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    
    left_content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(5))
    tf = left_content.text_frame
    tf.text = "User Pages:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "🏠 Home - Hero section & featured PGs"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "📋 Listings - Browse with filters"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "📄 PG Details - Full information"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "👤 Profile - User profile"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "🔐 Student Login"
    p.level = 1
    
    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    tf = right_content.text_frame
    tf.text = "Owner Pages:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = purple
    p = tf.add_paragraph()
    p.text = "📊 Dashboard - Manage listings"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "➕ Add PG - Create new listing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✏️ Edit PG - Update existing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "🔐 Owner Login"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Booking Confirmation"
    p.level = 1
    
    # Slide 9: Advanced Filtering System
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Advanced Filtering System"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    
    left_content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(5))
    tf = left_content.text_frame
    tf.text = "Filter Options:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "💰 Budget Range (Min/Max)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "📍 Location (City)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "👥 Gender (Boys/Girls/Unisex)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✨ Amenities (WiFi, AC, Food, etc.)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "🔍 Text Search (Name/Area)"
    p.level = 1
    
    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    tf = right_content.text_frame
    tf.text = "Sorting Options:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Price: Low to High"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Price: High to Low"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Rating: High to Low"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Rating: Low to High"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Note: Real-time Updates - Filters and sorting work together seamlessly"
    p.font.italic = True
    p.font.size = Pt(12)
    
    # Slide 10: State Management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "State Management with Context API"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "PGContext.jsx - Centralized State Management:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• listings: All PG data"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• loading: Loading state"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• addPG, updatePG, deletePG: CRUD operations"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• addReview, getReviews: Review management"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• addBooking, getBookings: Booking system"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• updateBookingStatus: Owner booking actions"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Benefits:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "• Global State: All components access PG data through context"
    p.level = 1
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "• No Prop Drilling: Clean component hierarchy"
    p.level = 1
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "• LocalStorage: Reviews and bookings persist in browser"
    p.level = 1
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "• Real-time Updates: Changes reflect immediately across app"
    p.level = 1
    p.font.name = "Arial"
    
    # Slide 11: User Journey
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "User Journey"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Browse: Visit home page → View featured PGs"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "2. Search: Go to listings → Apply filters → Sort results"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "3. Explore: Click PG card → View details, images, amenities"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "4. Review: Read reviews → Add your own review"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "5. Book: Login → Fill booking form → Get confirmation"
    p.font.size = Pt(18)
    
    # Slide 12: Owner Workflow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Owner Workflow"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Login: Owner login page → Access dashboard"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "2. Manage: View all listings → See statistics"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "3. Add/Edit: Create new PG or edit existing → Update details"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "4. Bookings: View booking requests → Approve/Reject"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "5. Status: Update PG status (Vacant/Occupied) → Manage availability"
    p.font.size = Pt(18)
    
    # Slide 13: UI/UX Design Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "UI/UX Design Highlights"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🎨 Modern Design: Clean, professional interface with purple/indigo theme"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "📱 Responsive: Works seamlessly on desktop, tablet, and mobile"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "⚡ Fast & Smooth: Optimized performance with Vite build tool"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "✨ Interactive: Hover effects, transitions, and smooth animations"
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Key Features: Gradient hero sections, card-based layouts, modal forms, status badges, and intuitive navigation"
    p.font.italic = True
    p.font.size = Pt(14)
    
    # Slide 14: Data Storage
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Data Storage & Management"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    
    left_content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(5))
    tf = left_content.text_frame
    tf.text = "PG Listings:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Stored in public/pgs.json"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• Loaded via fetch() on app start"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• CRUD operations update in-memory state"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• Can be extended to API backend"
    p.level = 1
    p.font.name = "Courier New"
    
    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5))
    tf = right_content.text_frame
    tf.text = "User Data:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "• Reviews: localStorage (per PG)"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• Bookings: localStorage"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• Login status: localStorage"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = "• Persists across sessions"
    p.level = 1
    p.font.name = "Courier New"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Note: Frontend-only application. Can easily be extended with backend API integration."
    p.font.italic = True
    p.font.size = Pt(11)
    p.font.name = "Arial"
    
    # Slide 15: Reusable Components
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Reusable Components"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "PGCard: Displays PG information, Image, name, location, Price and rating, Click to view details"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "FilterPanel: Budget range inputs, Location dropdown, Gender selection, Amenities checkboxes"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Navbar: Logo and navigation, Login buttons, User menu, Responsive design"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Icons: Custom SVG icons, Verified, Star, Search, Location, etc., Reusable across app"
    p.font.size = Pt(16)
    
    # Slide 16: Technical Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technical Implementation Highlights"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "React Hooks: useState, useEffect, useContext, useMemo for efficient state management"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "React Router: Dynamic routing with URL parameters (/pg/:id, /edit-pg/:id)"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Protected Routes: Dashboard redirects to login if not authenticated"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Memoization: useMemo for expensive filter/sort operations"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Form Handling: Controlled components with validation"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Error Handling: Image fallbacks, 404 pages, loading states"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "Responsive Design: TailwindCSS breakpoints (md:, lg:)"
    p.font.size = Pt(16)
    
    # Slide 17: Complete Feature List
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Complete Feature List"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    
    left_content = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4), Inches(5.5))
    tf = left_content.text_frame
    tf.text = "Search & Discovery:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "✅ Text search (name, location, area)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Multi-filter system"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Sort by price/rating"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Featured PGs on home"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Real-time result count"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "PG Management:"
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "✅ Add new PG listing"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Edit existing PG"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Delete PG"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Update status"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Multiple images"
    p.level = 1
    
    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4), Inches(5.5))
    tf = right_content.text_frame
    tf.text = "Reviews & Ratings:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "✅ Add reviews"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ View all reviews"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Star rating (1-5)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Average rating calculation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Review persistence"
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Booking System:"
    p.font.bold = True
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = "✅ Request booking"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Booking form"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Owner approval/rejection"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Booking confirmation"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "✅ Booking history"
    p.level = 1
    
    # Slide 18: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🔌 Backend Integration: REST API or GraphQL, Database (MongoDB/PostgreSQL), User authentication (JWT)"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "💳 Payment Gateway: Online payment integration, Booking deposits, Payment history"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "🗺️ Maps Integration: Google Maps API, Location visualization, Distance calculation"
    p.font.size = Pt(16)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "📧 Notifications: Email notifications, SMS alerts, In-app notifications"
    p.font.size = Pt(16)
    
    # Slide 19: Getting Started
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Getting Started"
    title.text_frame.paragraphs[0].font.color.rgb = purple
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Installation:"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(18)
    p = tf.add_paragraph()
    p.text = "npm install"
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.level = 1
    p = tf.add_paragraph()
    p.text = "npm run dev"
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.level = 1
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Usage:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "• Users: Browse → Search → Filter → View Details → Book"
    p.level = 1
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "• Owners: Login → Dashboard → Add/Edit PGs → Manage Bookings"
    p.level = 1
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Build for Production:"
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.name = "Arial"
    p = tf.add_paragraph()
    p.text = "npm run build"
    p.font.name = "Courier New"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 20: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    
    title_box = slide.shapes.add_textbox(left, top, width, Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Thank You!"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = purple
    title_para.alignment = PP_ALIGN.CENTER
    
    summary_box = slide.shapes.add_textbox(left, Inches(3.5), width, Inches(0.8))
    summary_frame = summary_box.text_frame
    summary_frame.text = "A complete, production-ready PG accommodation finder application"
    summary_para = summary_frame.paragraphs[0]
    summary_para.font.size = Pt(20)
    summary_para.font.color.rgb = dark_gray
    summary_para.alignment = PP_ALIGN.CENTER
    
    highlights_box = slide.shapes.add_textbox(left, Inches(4.5), width, Inches(2))
    highlights_frame = highlights_box.text_frame
    highlights_frame.text = "Built with: React, React Router, TailwindCSS"
    p = highlights_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = gray
    p.alignment = PP_ALIGN.CENTER
    p = highlights_frame.add_paragraph()
    p.text = "Features: Search, Filter, Book, Manage"
    p.font.size = Pt(16)
    p.font.color.rgb = gray
    p.alignment = PP_ALIGN.CENTER
    p = highlights_frame.add_paragraph()
    p.text = "Status: Frontend Complete, Ready for Backend"
    p.font.size = Pt(16)
    p.font.color.rgb = gray
    p.alignment = PP_ALIGN.CENTER
    
    questions_box = slide.shapes.add_textbox(left, Inches(6.5), width, Inches(0.5))
    questions_frame = questions_box.text_frame
    questions_frame.text = "Questions?"
    questions_para = questions_frame.paragraphs[0]
    questions_para.font.size = Pt(24)
    questions_para.font.color.rgb = gray
    questions_para.alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    print("Generating PowerPoint presentation...")
    prs = create_presentation()
    output_file = "PG_Finder_Presentation_With_Diagram.pptx"
    prs.save(output_file)
    print(f"Presentation saved as {output_file}")

